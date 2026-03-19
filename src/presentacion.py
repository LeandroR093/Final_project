from pptx import Presentation
from pptx.util import Inches, Pt

# 1. Crear una nueva presentación en blanco
prs = Presentation()

# --- DIAPOSITIVA 1: PORTADA ---
slide_1 = prs.slides.add_slide(prs.slide_layouts[0]) # Layout 0 es Portada
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]
title_1.text = "Oráculo Financiero AI 📈"
subtitle_1.text = "Predicción direccional del mercado de valores mediante Machine Learning\n\nAutores: Enrique Sucre y Leandro Roldan\nProyecto Final - 4Geeks Academy"

# --- DIAPOSITIVA 2: ÍNDICE ---
slide_2 = prs.slides.add_slide(prs.slide_layouts[1]) # Layout 1 es Título y Contenido
slide_2.shapes.title.text = "Índice del Proyecto"
tf_2 = slide_2.placeholders[1].text_frame
tf_2.text = "1. El Problema\n2. Los Datos\n3. Análisis Exploratorio (EDA)\n4. Combate de Modelos\n5. Optimización\n6. Despliegue en Vivo\n7. Conclusión"

# --- DIAPOSITIVA 3: EL PROBLEMA ---
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
slide_3.shapes.title.text = "1. El Problema a Resolver"
tf_3 = slide_3.placeholders[1].text_frame
tf_3.text = "Los mercados financieros son volátiles."
p = tf_3.add_paragraph()
p.text = "¿Puede la IA aprender de indicadores técnicos (RSI, MACD)?"
p = tf_3.add_paragraph()
p.text = "Objetivo: Predecir si la acción subirá 🟢 o bajará 🔴 mañana."

# --- DIAPOSITIVA 4: LOS DATOS ---
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
slide_4.shapes.title.text = "2. Los Datos y Preparación"
tf_4 = slide_4.placeholders[1].text_frame
tf_4.text = "Datos históricos de 500 empresas de Wall Street."
p = tf_4.add_paragraph()
p.text = "Label Encoding (pd.factorize): Transformamos nombres de texto (AAPL) a números (0, 1...) para que el modelo los entienda."

# --- DIAPOSITIVA 5: MODELOS ---
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
slide_5.shapes.title.text = "3. El Combate de Modelos"
tf_5 = slide_5.placeholders[1].text_frame
tf_5.text = "Enfrentamos a dos gigantes de los árboles de decisión:"
p = tf_5.add_paragraph()
p.text = "🌲 Random Forest: 72% de precisión (~79 segundos)"
p = tf_5.add_paragraph()
p.text = "🚀 XGBoost: 74% de precisión (~3.45 segundos) -> ¡GANADOR POR VELOCIDAD!"

# --- DIAPOSITIVA 6: OPTIMIZACIÓN ---
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
slide_6.shapes.title.text = "4. Optimización de Hiperparámetros"
tf_6 = slide_6.placeholders[1].text_frame
tf_6.text = "Usamos GridSearchCV para encontrar la receta perfecta."
p = tf_6.add_paragraph()
p.text = "Mejor configuración: learning_rate=0.1, max_depth=7, n_estimators=200."
p = tf_6.add_paragraph()
p.text = "Resultado Final: 75% de precisión (Accuracy)."

# --- DIAPOSITIVA 7: DESPLIEGUE ---
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
slide_7.shapes.title.text = "5. Despliegue y Conclusión"
tf_7 = slide_7.placeholders[1].text_frame
tf_7.text = "App interactiva construida con Streamlit."
p = tf_7.add_paragraph()
p.text = "Conclusión: El análisis técnico tiene poder predictivo a corto plazo usando XGBoost."
p = tf_7.add_paragraph()
p.text = "Disclaimer: Herramienta de apoyo, no recomendación de inversión."

# 2. Guardar el archivo físicamente en tu computadora
prs.save('Presentacion_Proyecto_Final.pptx')

print("✅ ¡Éxito! Tu archivo 'Presentacion_Proyecto_Final.pptx' ha sido creado.")